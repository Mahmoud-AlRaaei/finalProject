from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# إنشاء العرض التقديمي
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# تعريف الألوان
BLUE_DARK = RGBColor(0, 51, 102)      # أزرق داكن
BLUE_MAIN = RGBColor(0, 102, 204)     # أزرق رئيسي
BLUE_LIGHT = RGBColor(173, 216, 230)  # أزرق فاتح
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(64, 64, 64)
ACCENT = RGBColor(220, 20, 60)        # أحمر للتأكيد

def add_title_slide(prs, title, subtitle=""):
    """إضافة سلايد غلاف"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # الخلفية
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_DARK
    
    # العنوان
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # العنوان الفرعي
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = BLUE_LIGHT
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """إضافة سلايد محتوى"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # الخلفية
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # شريط العنوان
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                         Inches(0), Inches(0), 
                                         Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = BLUE_MAIN
    title_shape.line.color.rgb = BLUE_DARK
    
    # نص العنوان
    text_frame = title_shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT
    p.level = 0
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    
    # المحتوى
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(12)

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """إضافة سلايد بعمودين"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # الخلفية
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # شريط العنوان
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                         Inches(0), Inches(0), 
                                         Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = BLUE_MAIN
    title_shape.line.color.rgb = BLUE_DARK
    
    text_frame = title_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    
    # العمود الأيسر
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5.8))
    text_frame = left_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE_MAIN
    
    for item in left_content:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # العمود الأيمن
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    text_frame = right_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE_MAIN
    
    for item in right_content:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.space_before = Pt(8)
        p.space_after = Pt(8)

# ==================== SLIDE 1: الغلاف ====================
add_title_slide(prs, 
    "نظام ذكاء صنعي لاكتشاف الأخبار الكاذبة",
    "باللغة العربية على منصة X\n\nمشروع تخرج - كلية الهندسة المعلوماتية\nقسم الذكاء الاصطناعي - جامعة حمص\n\n2025-2026")

# ==================== SLIDE 2: فريق العمل ====================
add_content_slide(prs, "فريق العمل", [
    "👨‍💻 الطلاب:",
    "   • محمود الراعي",
    "   • عبد السلام إبراهيم",
    "   • بلال جوخدار",
    "   • علي العباس",
    "   • محمد الأبوحسنة",
    "",
    "👨‍🏫 الإشراف: د. ناصر أبو صالح"
])

# ==================== SLIDE 3: أهمية المشكلة ====================
add_content_slide(prs, "أهمية المشكلة", [
    "📊 الانتشار المتسارع للأخبار الكاذبة",
    "   • تأثر الرأي العام والمجتمع",
    "   • التضليل والحرب المعلوماتية",
    "",
    "🌍 التحديات في اللغة العربية",
    "   • نقص الأدوات والموارد المتخصصة",
    "   • تعقيد البنية اللغوية العربية",
    "",
    "🎯 الحاجة الماسة للحلول الذكية"
])

# ==================== SLIDE 4: أهداف المشروع ====================
add_content_slide(prs, "أهداف المشروع", [
    "✅ تطوير نموذج ذكاء صنعي متخصص للغة العربية",
    "",
    "✅ تحقيق دقة عالية في الكشف عن الأخبار الكاذبة",
    "",
    "✅ بناء تطبيق ويب سهل الاستخدام",
    "",
    "✅ توفير أداة عملية للمستخدمين والمحررين"
])

# ==================== SLIDE 5: مخطط العرض ====================
add_content_slide(prs, "مخطط العرض", [
    "1️⃣ الخلفية النظرية",
    "   • الأخبار الكاذبة والتحديات",
    "",
    "2️⃣ جمع وتحضير البيانات",
    "",
    "3️⃣ بناء وتدريب النموذج",
    "",
    "4️⃣ التطبيق العملي والواجهة",
    "",
    "5️⃣ النتائج والتقييم"
])

# ==================== SLIDE 6: الأخبار الكاذبة ====================
add_content_slide(prs, "ما هي الأخبار الكاذبة؟", [
    "تعريف:",
    "   معلومات مضللة أو غير دقيقة تنشر بقصد أو بدون قصد",
    "",
    "الأنواع:",
    "   • التضليل: معلومات كاذبة بقصد",
    "   • سوء المعلومات: معلومات خاطئة بدون قصد",
    "",
    "التأثير:",
    "   • التأثير على السياسات والقرارات",
    "   • تصعيد الصراعات والانقسامات"
])

# ==================== SLIDE 7: التحديات ====================
add_content_slide(prs, "تحديات الكشف", [
    "🔴 التحديات اللغوية:",
    "   • الكتابة العفوية والأخطاء الإملائية",
    "   • الكلمات المختلفة والمرادفات",
    "",
    "🟠 التحديات الفنية:",
    "   • نقص البيانات العربية المصنفة",
    "   • التشابه بين الأخبار الحقيقية والكاذبة",
    "",
    "🟡 السياق والتفاصيل الدقيقة"
])

# ==================== SLIDE 8: NLP والتعلم الآلي ====================
add_content_slide(prs, "معالجة اللغة الطبيعية (NLP)", [
    "تعريف:",
    "   فرع من الذكاء الاصطناعي يتعامل مع النصوص",
    "",
    "التطبيقات:",
    "   • تصنيف النصوص",
    "   • الترجمة الآلية",
    "   • تحليل المشاعر",
    "",
    "النموذج المستخدم: BERT و AraBERT"
])

# ==================== SLIDE 9: AraBERT ====================
add_content_slide(prs, "نموذج AraBERT", [
    "ما هو AraBERT؟",
    "   نموذج BERT متخصص للغة العربية",
    "",
    "خصائصه:",
    "   • تدريب على مليارات الكلمات العربية",
    "   • فهم عميق لسياق النصوص",
    "   • يعطي تمثيلات قوية للنصوص",
    "",
    "الميزة:",
    "   أداء ممتاز في المهام العربية"
])

# ==================== SLIDE 10: مصادر البيانات ====================
add_content_slide(prs, "جمع البيانات", [
    "مصادر البيانات:",
    "   • منصة X (تويتر)",
    "   • وسائل إعلام عربية",
    "   • مواقع أخبار موثوقة",
    "",
    "الحجم:",
    "   📊 6,000 خبر متوازن",
    "",
    "التصنيف:",
    "   ✓ 3,000 خبر صادق",
    "   ✗ 3,000 خبر كاذب"
])

# ==================== SLIDE 11: معالجة النصوص ====================
add_content_slide(prs, "معالجة النص العربي", [
    "خطوات التنظيف:",
    "   1. إزالة الروابط والهاشتاجات",
    "   2. إزالة الأرقام والرموز",
    "   3. توحيد الحروف (أ، إ، آ → ا)",
    "   4. إزالة التشكيل والحروف المشددة",
    "   5. إزالة المسافات الزائدة",
    "",
    "النتيجة: نصوص نظيفة وموحدة"
])

# ==================== SLIDE 12: توازن البيانات ====================
add_two_column_slide(prs, "توازن البيانات",
    "البيانات الأصلية", [
        "• الأخبار الصادقة: 50%",
        "• الأخبار الكاذبة: 50%",
        "",
        "• إجمالي: 6,000 خبر",
        "• متوازنة تماماً"
    ],
    "الفوائد", [
        "• تجنب التحيز",
        "• تدريب متوازن",
        "• نموذج عادل",
        "",
        "• أداء محسّن",
        "• نتائج موثوقة"
    ])

# ==================== SLIDE 13: معمارية النموذج ====================
add_content_slide(prs, "معمارية النموذج", [
    "المكونات:",
    "",
    "1️⃣ AraBERT Base Model",
    "   • طبقات تمثيل قوية",
    "",
    "2️⃣ Classification Head",
    "   • طبقة تصنيف ثنائي (صادق/كاذب)",
    "",
    "3️⃣ Softmax Output",
    "   • احتمالية الثقة للنتيجة"
])

# ==================== SLIDE 14: معاملات التدريب ====================
add_content_slide(prs, "معاملات التدريب (Hyperparameters)", [
    "⚙️ إعدادات التدريب:",
    "",
    "   • معدل التعلم: 2e-5",
    "   • حجم الدفعة: 16",
    "   • عدد الحقب: 10",
    "   • أقصى طول للنص: 128 رمز",
    "",
    "🔧 محسّن الوزن:",
    "   • AdamW Optimizer"
])

# ==================== SLIDE 15: عملية التدريب ====================
add_content_slide(prs, "عملية التدريب", [
    "📈 رحلة التدريب:",
    "",
    "   1. تحميل البيانات",
    "   2. تقسيم البيانات (80% تدريب، 20% اختبار)",
    "   3. حساب الخسارة والعودية",
    "   4. تحديث الأوزان",
    "   5. قياس الدقة على بيانات الاختبار",
    "",
    "⏱️ الوقت: حوالي 2-3 ساعات على GPU"
])

# ==================== SLIDE 16: النتائج المذهلة ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# شريط العنوان
title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     Inches(0), Inches(0), 
                                     Inches(10), Inches(0.9))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = BLUE_MAIN
text_frame = title_shape.text_frame
p = text_frame.paragraphs[0]
p.text = "النتائج النهائية"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.RIGHT
text_frame.margin_left = Inches(0.3)

# المستطيل الأكبر للنتيجة الرئيسية
result_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1.5), Inches(1.8),
                                    Inches(7), Inches(2.5))
result_box.fill.solid()
result_box.fill.fore_color.rgb = BLUE_LIGHT
result_box.line.color.rgb = BLUE_MAIN
result_box.line.width = Pt(3)

text_frame = result_box.text_frame
p = text_frame.paragraphs[0]
p.text = "دقة النموذج"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = BLUE_DARK
p.alignment = PP_ALIGN.CENTER

p = text_frame.add_paragraph()
p.text = "99.67%"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

# إحصائيات إضافية
stats_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(2))
text_frame = stats_box.text_frame
text_frame.word_wrap = True

stats = [
    "✓ الحساسية (Recall): 99.5%     |     ✓ التخصصية (Precision): 99.8%",
    "✓ F1-Score: 99.6%                 |     ✓ دقة الاختبار: 99.67%"
]

for stat in stats:
    p = text_frame.add_paragraph() if stats.index(stat) > 0 else text_frame.paragraphs[0]
    p.text = stat
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)

# ==================== SLIDE 17: الهندسة المعمارية للنظام ====================
add_content_slide(prs, "الهندسة المعمارية", [
    "🏗️ البنية الكاملة:",
    "",
    "Frontend (واجهة المستخدم)",
    "     ↓ HTML/CSS/JavaScript",
    "Flask Backend Server",
    "     ↓ معالجة الطلبات",
    "AraBERT Model",
    "     ↓ التنبؤ",
    "النتيجة (صادق/كاذب + نسبة الثقة)"
])

# ==================== SLIDE 18: تدفق البيانات ====================
add_content_slide(prs, "آلية عمل النظام", [
    "1️⃣ المستخدم يدخل النص",
    "",
    "2️⃣ معالجة النص",
    "   • تنظيف وتوحيد البيانات",
    "",
    "3️⃣ Tokenization",
    "   • تقسيم النص إلى رموز",
    "",
    "4️⃣ التنبؤ",
    "   • تمرير عبر النموذج",
    "",
    "5️⃣ النتيجة",
    "   • عرض التصنيف ونسبة الثقة"
])

# ==================== SLIDE 19: الميزات الرئيسية ====================
add_content_slide(prs, "الميزات الرئيسية", [
    "🎯 الكشف الفوري",
    "   • نتيجة خلال أقل من ثانية",
    "",
    "📊 نسبة الثقة",
    "   • إظهار مستوى التأكد من التصنيف",
    "",
    "🌐 واجهة سهلة",
    "   • تصميم استجابي جميل",
    "",
    "🔒 الخصوصية",
    "   • عدم حفظ أي بيانات"
])

# ==================== SLIDE 20: الواجهة الرسومية ====================
add_content_slide(prs, "واجهة المستخدم (UI)", [
    "📱 تصميم احترافي:",
    "",
    "✓ صندوق إدخال النص",
    "   اكتب أو الصق الخبر",
    "",
    "✓ زر التحليل",
    "   اضغط للحصول على النتيجة",
    "",
    "✓ عرض النتيجة",
    "   صادق ✓ أو كاذب ✗ مع نسبة الثقة"
])

# ==================== SLIDE 21: السرعة والأداء ====================
add_two_column_slide(prs, "السرعة والأداء",
    "على GPU", [
        "⚡ وقت الاستجابة:",
        "   < 100 ميلي ثانية",
        "",
        "🚀 معالجة متزامنة",
        "   تطبيق بدون تأخير",
        "",
        "🔄 عدد الطلبات",
        "   100+ طلب/دقيقة"
    ],
    "على CPU", [
        "⏱️ وقت الاستجابة:",
        "   500-800 ميلي ثانية",
        "",
        "💻 يعمل على أي جهاز",
        "   بدون احتياجات خاصة",
        "",
        "✅ نتائج موثوقة",
        "   دقة متطابقة"
    ])

# ==================== SLIDE 22: مقاييس الأداء ====================
add_content_slide(prs, "مقاييس التقييم", [
    "📊 Precision (التخصصية):",
    "   من الأخبار المتنبأ بها كـ 'كاذبة'، كم كانت صحيحة فعلاً؟",
    "   → 99.8%",
    "",
    "📊 Recall (الحساسية):",
    "   من الأخبار الكاذبة الفعلية، كم منها اكتشف النموذج؟",
    "   → 99.5%",
    "",
    "📊 F1-Score:",
    "   متوسط متوازن بين الاثنين → 99.6%"
])

# ==================== SLIDE 23: أمثلة عملية ====================
add_content_slide(prs, "أمثلة عملية", [
    "✅ مثال 1 (خبر صادق):",
    "   الإدخال: 'إطلاق البرنامج الوطني للذكاء الاصطناعي'",
    "   النتيجة: صادق (ثقة: 98.5%)",
    "",
    "❌ مثال 2 (خبر كاذب):",
    "   الإدخال: 'اكتشاف طريقة سحرية للثراء'",
    "   النتيجة: كاذب (ثقة: 99.2%)",
    "",
    "ملاحظة: الأمثلة توضيحية"
])

# ==================== SLIDE 24: المقارنة مع الطرق السابقة ====================
add_two_column_slide(prs, "المقارنة",
    "الطرق التقليدية", [
        "• الكشف اليدوي",
        "   بطيء وغير فعال",
        "",
        "• القوائم السوداء",
        "   غير شاملة",
        "",
        "• النماذج البسيطة",
        "   دقة حوالي 75-85%"
    ],
    "نموذجنا (AraBERT)", [
        "• كشف آلي فوري",
        "   سريع وفعال",
        "",
        "• يتعلم من البيانات",
        "   قابل للتطوير",
        "",
        "• نموذج متقدم",
        "   دقة 99.67%"
    ])

# ==================== SLIDE 25: التحديات والحلول ====================
add_two_column_slide(prs, "التحديات والحلول",
    "التحديات", [
        "🔴 نقص البيانات العربية",
        "",
        "🔴 تعقيد اللغة العربية",
        "",
        "🔴 التطور السريع للأخبار",
        "",
        "🔴 الموارد الحاسوبية"
    ],
    "الحلول المطبقة", [
        "✓ جمع 6,000 خبر متوازن",
        "",
        "✓ معالجة نصوص متقدمة",
        "",
        "✓ نموذج قابل للتحديث",
        "",
        "✓ دعم GPU و CPU"
    ])

# ==================== SLIDE 26: الإنجازات ====================
add_content_slide(prs, "الإنجازات الرئيسية", [
    "🏆 إنجازات المشروع:",
    "",
    "✅ نموذج دقته 99.67%",
    "",
    "✅ نظام كامل متكامل",
    "",
    "✅ واجهة ويب احترافية",
    "",
    "✅ توثيق شامل",
    "",
    "✅ قابل للتطوير والتحسين"
])

# ==================== SLIDE 27: التطويرات المستقبلية ====================
add_content_slide(prs, "التطويرات المستقبلية", [
    "🚀 خطط التطوير:",
    "",
    "📱 تطبيق موبايل",
    "",
    "🔄 تحديث النموذج بانتظام",
    "",
    "🌐 دعم لغات أخرى",
    "",
    "📊 لوحة تحكم متقدمة",
    "",
    "🤝 تكامل مع منصات إخبارية"
])

# ==================== SLIDE 28: الشكر والختام ====================
add_title_slide(prs,
    "شكراً لاهتمامكم",
    "نظام ذكاء صنعي لاكتشاف الأخبار الكاذبة باللغة العربية\n\nللأسئلة والاستفسارات\n\nشكراً للدكتور ناصر أبو صالح على الإشراف والتوجيه")

# ==================== حفظ الملف ====================
prs.save('presentation.pptx')
print("✅ تم إنشاء العرض التقديمي بنجاح!")
print("📁 اسم الملف: presentation.pptx")
print("📊 عدد السلايدات: 28")
print("🎨 الألوان: أزرق احترافي")
