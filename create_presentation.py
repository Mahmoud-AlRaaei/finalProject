from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

# إنشاء عرض احترافي
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# تعريف الألوان الاحترافية
BLUE_DARK = RGBColor(0, 51, 102)       # أزرق داكن جداً
BLUE_MAIN = RGBColor(0, 102, 204)      # أزرق رئيسي
BLUE_LIGHT = RGBColor(230, 242, 255)   # أزرق فاتح جداً
BLUE_ACCENT = RGBColor(0, 153, 255)    # أزرق زاهي
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(45, 45, 45)
MEDIUM_GRAY = RGBColor(100, 100, 100)
LIGHT_GRAY = RGBColor(240, 240, 240)
SUCCESS_GREEN = RGBColor(46, 184, 92)
WARNING_ORANGE = RGBColor(243, 156, 18)
DANGER_RED = RGBColor(231, 76, 60)

def add_background(slide, color=WHITE):
    """إضافة خلفية للسلايد"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title, color=BLUE_MAIN):
    """إضافة شريط عنوان احترافي"""
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = color
    title_shape.line.color.rgb = color
    title_shape.line.width = Pt(0)
    
    text_frame = title_shape.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT
    text_frame.margin_left = Inches(0.5)
    text_frame.margin_right = Inches(0.5)

def add_content_bullet(slide, title, bullets):
    """إضافة سلايد محتوى مع نقاط"""
    add_background(slide, LIGHT_GRAY)
    add_title_bar(slide, title)
    
    # منطقة المحتوى
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.3),
        Inches(8.6), Inches(5.8)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(10)
        p.space_after = Pt(10)

def add_two_col_slides(slide, title, left_title, left_items, right_title, right_items):
    """إضافة سلايد بعمودين احترافي"""
    add_background(slide, LIGHT_GRAY)
    add_title_bar(slide, title)
    
    # العمود الأيسر
    left_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.3),
        Inches(4.5), Inches(5.8)
    )
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = WHITE
    left_shape.line.color.rgb = BLUE_ACCENT
    left_shape.line.width = Pt(3)
    
    left_text = left_shape.text_frame
    left_text.word_wrap = True
    left_text.margin_left = Inches(0.2)
    left_text.margin_right = Inches(0.2)
    
    p = left_text.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE_MAIN
    p.space_after = Pt(15)
    
    for item in left_items:
        p = left_text.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(5)
        p.space_after = Pt(5)
    
    # العمود الأيمن
    right_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(1.3),
        Inches(4.5), Inches(5.8)
    )
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = WHITE
    right_shape.line.color.rgb = BLUE_ACCENT
    right_shape.line.width = Pt(3)
    
    right_text = right_shape.text_frame
    right_text.word_wrap = True
    right_text.margin_left = Inches(0.2)
    right_text.margin_right = Inches(0.2)
    
    p = right_text.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE_MAIN
    p.space_after = Pt(15)
    
    for item in right_items:
        p = right_text.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(5)
        p.space_after = Pt(5)

def add_stat_slide(slide, title, stat_value, stat_label, description):
    """إضافة سلايد إحصائية كبيرة"""
    add_background(slide, LIGHT_GRAY)
    add_title_bar(slide, title, BLUE_DARK)
    
    # مربع الإحصائية الكبير
    stat_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(8), Inches(3.5)
    )
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = WHITE
    stat_box.line.color.rgb = BLUE_ACCENT
    stat_box.line.width = Pt(4)
    
    text_frame = stat_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = stat_value
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = stat_label
    p.font.size = Pt(32)
    p.font.color.rgb = BLUE_MAIN
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)
    
    # الوصف تحت
    desc_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.2),
        Inches(8), Inches(1.8)
    )
    desc_text = desc_box.text_frame
    desc_text.word_wrap = True
    p = desc_text.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 1: غلاف مميز ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, BLUE_DARK)

# العنوان الرئيسي
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "نظام ذكاء صنعي لاكتشاف الأخبار الكاذبة"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "باللغة العربية على منصة X"
p.font.size = Pt(32)
p.font.color.rgb = BLUE_ACCENT
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(15)

# المعلومات السفلى
info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
tf = info_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "مشروع تخرج - كلية الهندسة المعلوماتية\nقسم الذكاء الاصطناعي - جامعة حمص"
p.font.size = Pt(22)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 2: الفريق ====================
add_content_bullet(prs.slides.add_slide(prs.slide_layouts[6]), "فريق العمل", [
    "👨‍💻 الطلاب المشاركون:",
    "   • محمود الراعي (المشرف التقني)",
    "   • عبد السلام إبراهيم",
    "   • بلال جوخدار",
    "   • علي العباس",
    "   • محمد الأبوحسنة",
    "",
    "👨‍🏫 الإشراف الأكاديمي:",
    "   د. ناصر أبو صالح",
    "   العام الدراسي 2025-2026"
])

# ==================== SLIDE 3: المشكلة ====================
add_content_bullet(prs.slides.add_slide(prs.slide_layouts[6]), "المشكلة والتحفيز", [
    "🔴 أزمة الأخبار الكاذبة:",
    "   • انتشار سريع على وسائل التواصل الاجتماعي",
    "   • تأثيرات سلبية على المجتمع والرأي العام",
    "",
    "⚠️ التحديات الخاصة باللغة العربية:",
    "   • نقص الأدوات والموارد المتخصصة",
    "   • تعقيد البنية اللغوية والصيغ المختلفة",
    "   • عدم وجود مجموعات بيانات عربية كبيرة",
    "",
    "✨ الفرصة:",
    "   استخدام تقنيات التعلم العميق الحديثة"
])

# ==================== SLIDE 4: الأهداف ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_GRAY)
add_title_bar(slide, "أهداف المشروع")

goals = [
    ("تطوير نموذج متخصص", "نموذج BERT مدرب على النصوص العربية"),
    ("دقة عالية", "الوصول إلى دقة 95%+ في الكشف"),
    ("سهولة الاستخدام", "واجهة ويب سهلة وسريعة"),
    ("قابلية التطوير", "نموذج قابل للتحديث والتحسين")
]

y_pos = 1.3
for goal_title, goal_desc in goals:
    # مربع الهدف
    goal_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(y_pos),
        Inches(8.6), Inches(1)
    )
    goal_box.fill.solid()
    goal_box.fill.fore_color.rgb = WHITE
    goal_box.line.color.rgb = BLUE_ACCENT
    goal_box.line.width = Pt(2)
    
    text_frame = goal_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = "✓ " + goal_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    
    p = text_frame.add_paragraph()
    p.text = goal_desc
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(5)
    
    y_pos += 1.15

# ==================== SLIDE 5: AraBERT ====================
add_content_bullet(prs.slides.add_slide(prs.slide_layouts[6]), "نموذج AraBERT", [
    "🤖 ما هو BERT؟",
    "   نموذج Bidirectional Encoder Representations from Transformers",
    "   مدرب على ملايين النصوص لفهم اللغة بعمق",
    "",
    "🌍 ما هو AraBERT؟",
    "   نسخة متخصصة من BERT للغة العربية",
    "   تم تدريبه على 66 مليار رمز عربي",
    "",
    "💡 المميزات:",
    "   • فهم سياق النصوص العربية",
    "   • تمثيلات قوية للكلمات والجمل",
    "   • معالجة التشكيل والصيغ المختلفة",
    "   • أداء ممتاز على المهام العربية"
])

# ==================== SLIDE 6: البيانات ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_GRAY)
add_title_bar(slide, "جمع وتحضير البيانات")

# إحصائيات البيانات
stats_data = [
    ("6,000", "إجمالي الأخبار", BLUE_MAIN),
    ("3,000", "أخبار صادقة ✓", SUCCESS_GREEN),
    ("3,000", "أخبار كاذبة ✗", DANGER_RED)
]

x_positions = [0.7, 3.65, 6.6]
for idx, (number, label, color) in enumerate(stats_data):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_positions[idx]), Inches(1.5),
        Inches(2.5), Inches(2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.width = Pt(0)
    
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# المصادر
sources_box = slide.shapes.add_textbox(Inches(0.7), Inches(4), Inches(8.6), Inches(3))
tf = sources_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📊 مصادر البيانات:"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = BLUE_MAIN
p.space_after = Pt(10)

sources = [
    "✓ منصة X (تويتر) - التغريدات والمناقشات",
    "✓ وسائل إعلام عربية موثوقة",
    "✓ مواقع أخبار متخصصة",
    "✓ جهود يدوية لتصنيف وتوثيق الأخبار"
]

for source in sources:
    p = tf.add_paragraph()
    p.text = source
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(5)

# ==================== SLIDE 7: معالجة النص ====================
add_content_bullet(prs.slides.add_slide(prs.slide_layouts[6]), "معالجة النص العربي (NLP)", [
    "🔧 خطوات التنظيف:",
    "",
    "1️⃣ إزالة الروابط والهاشتاجات والإشارات",
    "",
    "2️⃣ إزالة الأرقام والرموز الخاصة",
    "",
    "3️⃣ توحيد الحروف (أ، إ، آ → ا)",
    "",
    "4️⃣ إزالة التشكيل (الفتحة، الكسرة، الضمة)",
    "",
    "5️⃣ حذف الكلمات المكررة المتتالية",
    "",
    "6️⃣ تطبيع المسافات"
])

# ==================== SLIDE 8: معمارية النموذج ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_GRAY)
add_title_bar(slide, "معمارية النموذج")

# الرسم التخطيطي
layers = [
    ("النص العربي", BLUE_LIGHT, 1.3),
    ("Tokenization", BLUE_ACCENT, 2.5),
    ("AraBERT Encoder", BLUE_MAIN, 3.7),
    ("Classification Head", BLUE_DARK, 4.9),
    ("النتيجة (صادق/كاذب)", SUCCESS_GREEN, 6.1)
]

for text, color, y in layers:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(y),
        Inches(4), Inches(0.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.width = Pt(0)
    
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # السهم
    if y < 6.1:
        arrow = slide.shapes.add_connector(1, Inches(5), Inches(y + 0.8), Inches(5), Inches(y + 1.1))
        arrow.line.color.rgb = DARK_GRAY
        arrow.line.width = Pt(2)

# ==================== SLIDE 9: التدريب ====================
add_content_bullet(prs.slides.add_slide(prs.slide_layouts[6]), "معاملات التدريب (Hyperparameters)", [
    "⚙️ الإعدادات الرئيسية:",
    "",
    "   • معدل التعلم (Learning Rate): 2e-5",
    "   • حجم الدفعة (Batch Size): 16",
    "   • عدد الحقب (Epochs): 10",
    "   • أقصى طول للنص: 128 رمز",
    "   • محسّن الأوزان: AdamW",
    "",
    "📊 البيانات:",
    "   • 80% للتدريب (4,800 خبر)",
    "   • 20% للاختبار (1,200 خبر)",
    "",
    "⏱️ وقت التدريب: 2-3 ساعات على GPU"
])

# ==================== SLIDE 10: النتائج المذهلة ====================
add_stat_slide(prs.slides.add_slide(prs.slide_layouts[6]),
    "دقة النموذج النهائية",
    "99.67%",
    "دقة على بيانات الاختبار",
    "نتيجة استثنائية تعكس جودة التدريب والبيانات والمعمارية"
)

# ==================== SLIDE 11: مقاييس الأداء ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_GRAY)
add_title_bar(slide, "مقاييس التقييم الشاملة")

metrics = [
    ("Accuracy", "99.67%", "النسبة الإجمالية للتصنيفات الصحيحة"),
    ("Precision", "99.8%", "من الأخبار المتنبأ بها كـ 'كاذبة'، كم فعلاً كاذبة"),
    ("Recall", "99.5%", "من الأخبار الكاذبة الفعلية، كم منها اكتشف النموذج"),
    ("F1-Score", "99.6%", "متوسط متوازن بين Precision و Recall")
]

y_pos = 1.3
for metric_name, metric_value, description in metrics:
    # الخلفية
    bg_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(y_pos),
        Inches(8.6), Inches(1.2)
    )
    bg_box.fill.solid()
    bg_box.fill.fore_color.rgb = WHITE
    bg_box.line.color.rgb = BLUE_LIGHT
    bg_box.line.width = Pt(2)
    
    tf = bg_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = metric_name + " : " + metric_value
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    
    p = tf.add_paragraph()
    p.text = description
    p.font.size = Pt(13)
    p.font.color.rgb = MEDIUM_GRAY
    p.space_before = Pt(3)
    
    y_pos += 1.35

# ==================== SLIDE 12: الواجهة الرسومية ====================
add_content_bullet(prs.slides.add_slide(prs.slide_layouts[6]), "واجهة المستخدم (UI/UX)", [
    "🎨 تصميم استجابي وجميل:",
    "",
    "📝 صندوق إدخال النص:",
    "   • إمكانية كتابة أو لصق الخبر",
    "   • دعم النصوص العربية بالكامل",
    "",
    "🔍 زر التحليل:",
    "   • معالجة فورية",
    "   • رسالات تحميل جميلة",
    "",
    "📊 عرض النتيجة:",
    "   • صادق ✓ أو كاذب ✗ مع أيقونات واضحة",
    "   • نسبة الثقة (0-100%)",
    "   • تفسير بسيط للنتيجة"
])

# ==================== SLIDE 13: السرعة والأداء ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_GRAY)
add_title_bar(slide, "السرعة والأداء", BLUE_DARK)

two_col = [
    ("على GPU (معالج رسومي)", [
        "⚡ < 100 ميلي ثانية",
        "   وقت الاستجابة",
        "",
        "🚀 معالجة فورية",
        "   بدون تأخير ملحوظ",
        "",
        "📈 100+ طلب/دقيقة",
        "   معالجة متزامنة"
    ], "على CPU (معالج مركزي)", [
        "⏱️ 500-800 ميلي ثانية",
        "   وقت الاستجابة",
        "",
        "💻 يعمل على أي جهاز",
        "   بدون احتياجات خاصة",
        "",
        "✅ نتائج موثوقة",
        "   دقة متطابقة"
    ])

# تطبيق العمودين
add_two_col_slides(slide, "", two_col[0], two_col[1], two_col[2], two_col[3])

# ==================== SLIDE 14: آلية العمل ====================
add_content_bullet(prs.slides.add_slide(prs.slide_layouts[6]), "تدفق معالجة الأخبار", [
    "🔄 خطوات المعالجة:",
    "",
    "1️⃣ المستخدم يدخل النص",
    "",
    "2️⃣ تنظيف وتطبيع النص",
    "   إزالة الرموز والروابط والأرقام",
    "",
    "3️⃣ Tokenization",
    "   تقسيم النص إلى رموز (tokens)",
    "",
    "4️⃣ معالجة بـ AraBERT",
    "   استخراج خصائص النص (embeddings)",
    "",
    "5️⃣ التنبؤ",
    "   تمرير الخصائص للطبقة التصنيف",
    "",
    "6️⃣ النتيجة النهائية",
    "   صادق ✓ أو كاذب ✗ + نسبة الثقة"
])

# ==================== SLIDE 15: المميزات الرئيسية ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_GRAY)
add_title_bar(slide, "المميزات الرئيسية")

features = [
    "🎯 الكشف الفوري والدقيق",
    "📊 عرض نسبة الثقة والموثوقية",
    "🌐 واجهة ويب حديثة وسهلة",
    "🔒 عدم حفظ البيانات الشخصية",
    "⚡ سرعة معالجة عالية",
    "📱 متوافق مع جميع الأجهزة"
]

for idx, feature in enumerate(features):
    row = idx // 3
    col = idx % 3
    x_pos = 0.7 + (col * 3)
    y_pos = 1.5 + (row * 2.5)
    
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_pos), Inches(y_pos),
        Inches(2.8), Inches(2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = BLUE_ACCENT
    box.line.width = Pt(2)
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = tf.paragraphs[0]
    p.text = feature
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLUE_MAIN
    p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 16: التحديات والحلول ====================
add_two_col_slides(prs.slides.add_slide(prs.slide_layouts[6]),
    "التحديات والحلول",
    "التحديات 🔴", [
        "نقص البيانات العربية المصنفة",
        "",
        "تعقيد البنية اللغوية",
        "",
        "التطور السريع للأخبار",
        "",
        "الموارد الحاسوبية",
        "",
        "التشابه بين الحقيقي والكاذب"
    ],
    "الحلول ✅", [
        "جمع وتصنيف 6,000 خبر متوازن",
        "",
        "معالجة نصوص متقدمة ومتخصصة",
        "",
        "نموذج قابل للتحديث المستمر",
        "",
        "دعم GPU و CPU",
        "",
        "استخدام AraBERT المتقدم"
    ]
)

# ==================== SLIDE 17: الإنجازات ====================
add_content_bullet(prs.slides.add_slide(prs.slide_layouts[6]), "الإنجازات الرئيسية 🏆", [
    "✅ نموذج تعلم عميق متقدم",
    "   دقة 99.67% - نتيجة استثنائية",
    "",
    "✅ نظام متكامل وشامل",
    "   من البيانات إلى الإنتاج",
    "",
    "✅ واجهة ويب احترافية",
    "   تجربة مستخدم ممتازة",
    "",
    "✅ توثيق شامل وكامل",
    "   سهولة الصيانة والتطوير",
    "",
    "✅ قابل للتطوير والتحسين",
    "   أساس قوي للمستقبل"
])

# ==================== SLIDE 18: التطويرات المستقبلية ====================
add_content_bullet(prs.slides.add_slide(prs.slide_layouts[6]), "التطويرات المستقبلية 🚀", [
    "📱 تطبيق موبايل",
    "   iOS و Android تطبيقات",
    "",
    "🔄 تحديث النموذج بانتظام",
    "   التعلم من البيانات الجديدة",
    "",
    "🌐 دعم لغات أخرى",
    "   إنجليزي وفرنسي وغيرها",
    "",
    "📊 لوحة تحكم متقدمة",
    "   إحصائيات وتقارير مفصلة",
    "",
    "🤝 تكامل مع منصات إخبارية",
    "   API للاستخدام المباشر"
])

# ==================== SLIDE 19: الخاتمة ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, BLUE_DARK)

# النص الختامي
closing_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
tf = closing_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "شكراً لاهتمامكم"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = BLUE_ACCENT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "نظام ذكاء صنعي لاكتشاف الأخبار الكاذبة باللغة العربية"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)

p = tf.add_paragraph()
p.text = "للأسئلة والاستفسارات\nشكراً لـ د. ناصر أبو صالح على التوجيه والإشراف"
p.font.size = Pt(18)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(30)

# ==================== حفظ الملف ====================
prs.save('presentation_professional.pptx')
print("✅ تم إنشاء العرض التقديمي الاحترافي!")
print("📁 الملف: presentation_professional.pptx")
print("📊 عدد السلايدات: 19 سلايد شامل")
print("🎨 التصميم: احترافي وعصري")
print("💙 الألوان: أزرق احترافي مع إحصائيات ملونة")
